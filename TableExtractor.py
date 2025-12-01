# Import necessary libraries for Flask app and file handling
from flask import Flask, request, render_template_string, jsonify, flash, redirect
from werkzeug.utils import secure_filename
import os
import re
import docx
import pandas as pd
import pdfplumber
import json
from PIL import Image
import io
import pptx
import pythoncom
import comtypes.client
from bs4 import BeautifulSoup
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from dotenv import load_dotenv
import os

# Load environment variables from .env file
load_dotenv()  # ✅ Must be called first, before using os.getenv

# ✅ Now safely load the values
AZURE_ENDPOINT = os.getenv("AZURE_FORMRECOGNIZER_ENDPOINT")
AZURE_KEY = os.getenv("AZURE_FORMRECOGNIZER_KEY")

# Initialize Azure Form Recognizer client
azure_client = DocumentAnalysisClient(
    endpoint=AZURE_ENDPOINT,
    credential=AzureKeyCredential(AZURE_KEY)
)

# ✅ Debug print to verify
print("🔍 AZURE_ENDPOINT from env:", AZURE_ENDPOINT)
print("🔍 AZURE_KEY from env:", AZURE_KEY[:5], "***")

# Donut model imports for vision-based table extraction
from transformers import DonutProcessor, VisionEncoderDecoderModel

# Flask app initialization
app = Flask(__name__)
app.secret_key = 'your-secret-key-here'
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

# Supported file extensions for upload
ALLOWED_EXTENSIONS = {
    'pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'csv',
    'txt', 'html', 'htm', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp'
}

def allowed_file(filename):
    """Check if uploaded file is allowed based on its extension."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# ----------- DONUT MODEL LOADING (singleton) -----------
donut_processor = None
donut_model = None

def get_donut_model():
    """Load Donut model and processor only once (singleton)."""
    global donut_processor, donut_model
    if donut_processor is None or donut_model is None:
        donut_processor = DonutProcessor.from_pretrained("naver-clova-ix/donut-base-finetuned-docvqa")
        donut_model = VisionEncoderDecoderModel.from_pretrained("naver-clova-ix/donut-base-finetuned-docvqa")
    return donut_processor, donut_model

def get_donut_summary_text(file_path):
    """Use Donut to summarize document image."""
    try:
        processor, model = get_donut_model()
        image = Image.open(file_path).convert("RGB")
        prompt = "<s_docvqa><s_question>Read and summarize the document including field names, values, and table labels</s_question><s_answer>"
        inputs = processor(image, prompt, return_tensors="pt")
        outputs = model.generate(**inputs)
        result = processor.batch_decode(outputs, skip_special_tokens=True)[0]
        return result.strip()
    except Exception as e:
        print(f"Donut text summary error: {e}")
        return ""

def extract_tables_from_image_with_donut(file_path):
    """Extract tables from image using Donut model."""
    try:
        processor, model = get_donut_model()
        image = Image.open(file_path).convert("RGB")
        prompt = "<s_docvqa><s_question>Extract all table data with merged cells as JSON. Use headers, values. Use 'colspan' and 'rowspan' if needed. Output only JSON.</s_question><s_answer>"
        inputs = processor(image, prompt, return_tensors="pt")
        outputs = model.generate(**inputs)
        result = processor.batch_decode(outputs, skip_special_tokens=True)[0]
        # Try to get first JSON in output
        match = re.search(r'\{[\s\S]*\}', result)
        if not match:
            return []
        json_str = match.group(0)
        result_data = json.loads(json_str)
        headers = result_data.get("headers", [])
        data = result_data.get("data", [])
        merged = data if data and isinstance(data[0][0], dict) else None
        return [{
            'source': 'Image Table (Donut)',
            'headers': headers,
            'data': None,
            'merged': merged,
            'method': 'Donut'
        }]
    except Exception as e:
        print(f"Donut fallback error: {e}")
        return []

def extract_tables_from_image_with_llm(file_path):
    """
    Use only Azure for table extraction.
    Then pass Azure output through Gemini (or similar LLM) for output formatting if needed.
    """
    if not file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.webp', '.tiff', '.bmp')):
        print(f"⚠️ Skipping LLM processing for non-image file: {file_path}")
        return []

    # Step 1: Use Azure for table extraction
    azure_tables = extract_tables_with_azure_layout(file_path)
    if not azure_tables:
        print("⚠️ Azure returned no tables")
        return []

    # Step 2: Apply formatting using Gemini/LLM (optional, based on user instruction)
    # Here we simulate a formatting step (you can replace this with your own formatter logic)
    for table in azure_tables:
        for row in table.get("merged", []):
            for cell in row:
                val = cell.get("value", "")
                val = val.replace("* ", "").replace("*", "")
                val = val.replace("®", "")
                val = val.replace("¢", "")
                cell["value"] = val.strip()

    return azure_tables

def extract_tables_from_image_with_ocr(file_path):
    """OCR fallback for images to extract tables as text."""
    try:
        import pytesseract
        image = Image.open(file_path).convert("RGB")
        text = pytesseract.image_to_string(image)
        lines = text.split("\n")
        rows = [re.split(r'\s{2,}', line.strip()) for line in lines if line.strip()]
        if len(rows) < 2:
            return []
        headers = rows[0]
        data = [dict(zip(headers, row + [""] * (len(headers)-len(row)))) for row in rows[1:]]
        return [{
            'source': 'Image Table (OCR Fallback)',
            'headers': headers,
            'data': data,
            'merged': None,
            'method': 'OCR'
        }]
    except Exception as e:
        print(f"OCR fallback error: {e}")
        return []

def tables_are_incomplete(tables):
    """Check if extracted tables are incomplete."""
    if not tables:
        return True
    for t in tables:
        if not (t.get("data") or t.get("merged")) or not t.get("headers"):
            return True
        if t["headers"] == ["Output"]:
            return True
    return False


def convert_pptx_to_pdf(input_path, output_path):
    """Convert PPTX file to PDF using COM automation (Windows only)."""
    import os
    import comtypes.client
    import pythoncom

    pythoncom.CoInitialize()  # ✅ Place this inside the function

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)

    try:
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        presentation.SaveAs(output_path, 32)  # 32 = PDF format
        presentation.Close()
    finally:
        powerpoint.Quit()

def extract_tables_from_pptx_with_gemini(file_path):
    """Extract tables from PPTX using Gemini fallback (converts to PDF first)."""
    tables = []
    try:
        import comtypes.client
        # Convert PPTX to PDF, then each PDF page to image and run Gemini
        pdf_out = file_path + ".pdf"
        convert_pptx_to_pdf(file_path, pdf_out)
        pdf_tables = extract_tables_from_pdf_with_azure(pdf_out)
        for t in pdf_tables:
            t['source'] = f'PPTX->PDF Azure Fallback - {t["source"]}'
            t['method'] = "pptx2pdf+Azure"
        tables += pdf_tables
        os.remove(pdf_out)
    except Exception as e:
        print(f"PPTX Azure fallback error: {e}")
    return tables

from pdf2image import convert_from_path

def extract_tables_from_pdf(file_path):
    """Extract tables from PDF using pdfplumber. If none, convert each page to image and run Azure per page."""
    tables = []

    # Try primary method
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_tables = page.extract_tables()
                for table_num, table in enumerate(page_tables, 1):
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0] if table[0] else None)
                        df = df.fillna('')
                        tables.append({
                            'source': f'Page {page_num}, Table {table_num}',
                            'data': df.to_dict('records'),
                            'headers': df.columns.tolist(),
                            'merged': None,
                            'method': 'pdfplumber'
                        })
    except Exception as e:
        print(f"PDF extraction error with pdfplumber: {e}")

    # If no tables found, fallback: convert each PDF page to image and use Azure Layout
    if not tables:
        try:
            print("⚠️ pdfplumber found no tables. Falling back to per-page Azure Layout...")
            poppler_path = r"C:\Users\rohit\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin"  # ✅ Update to your poppler path
            page_images = convert_from_path(file_path, poppler_path=poppler_path)

            for idx, image in enumerate(page_images, 1):
                img_path = f"{file_path}.page{idx}.png"
                image.save(img_path)
                page_tables = extract_tables_with_azure_layout(img_path)
                for t in page_tables:
                    t["source"] = f"Azure Page {idx} - {t['source']}"
                    t["method"] = "Azure Layout (Per Page)"
                tables += page_tables
                os.remove(img_path)
        except Exception as e:
            print(f"Azure fallback per-page error: {e}")

    return tables



def extract_tables_from_pptx(file_path):
    """Extract tables from PPTX using python-pptx, fallback to Azure Layout if none found."""
    tables = []
    try:
        prs = pptx.Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape_num, shape in enumerate(slide.shapes, 1):
                if not shape.has_table:
                    continue
                table = shape.table
                rows = list(table.rows)
                if not rows:
                    continue
                header_row_idx = None
                for idx, row in enumerate(rows):
                    if any(cell.text.strip() for cell in row.cells):
                        header_row_idx = idx
                        break
                if header_row_idx is None:
                    continue
                headers = [cell.text.strip() or f"Column {i+1}" for i, cell in enumerate(rows[header_row_idx].cells)]
                data = []
                for row in rows[header_row_idx+1:]:
                    row_data = [cell.text.strip() for cell in row.cells]
                    while len(row_data) < len(headers):
                        row_data.append('')
                    while len(row_data) > len(headers):
                        headers.append(f"Column {len(headers)+1}")
                    data.append(dict(zip(headers, row_data)))
                tables.append({
                    'source': f'Slide {slide_num}, Shape {shape_num}',
                    'data': data,
                    'headers': headers,
                    'merged': None,
                    'method': 'python-pptx'
                })
    except Exception as e:
        print(f"Error extracting from PPTX using python-pptx: {e}")

    if not tables:
        print("⚠️ python-pptx found no tables. Falling back to Azure Layout...")
        try:
            tables = extract_tables_with_azure_layout(file_path)
            for t in tables:
                t['source'] = f'Azure Layout Fallback - {t["source"]}'
                t['method'] = 'Azure Layout'
        except Exception as e:
            print(f"Azure fallback error: {e}")

    return tables


def extract_tables_from_docx(file_path):
    """Extract tables from DOCX, fallback to PDF/image if none found."""
    tables = []
    try:
        doc = docx.Document(file_path)
        for table_num, table in enumerate(doc.tables, 1):
            data = []
            headers = []
            if table.rows:
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                for row in table.rows[1:]:
                    row_data = [cell.text.strip() for cell in row.cells]
                    data.append(dict(zip(headers, row_data)))
                tables.append({
                    'source': f'Table {table_num}',
                    'data': data,
                    'headers': headers,
                    'merged': None,
                    'method': 'python-docx'
                })
    except Exception as e:
        print(f"Error extracting from DOCX: {e}")
    if not tables:
        # fallback: convert docx to PDF/image and use PDF/image fallback
        try:
            from docx2pdf import convert
            pdf_out = file_path + ".pdf"
            convert(file_path, pdf_out)
            pdf_tables = extract_tables_from_pdf(pdf_out)
            for t in pdf_tables:
                t['source'] = f'DOCX->PDF Fallback - {t["source"]}'
                t['method'] = "docx2pdf+pdfplumber"
            tables += pdf_tables
            os.remove(pdf_out)
        except Exception as e:
            print(f"DOCX fallback error: {e}")
    return tables

def extract_tables_from_excel(file_path):
    """Extract tables from Excel, fallback to HTML if none found."""
    tables = []
    try:
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            if not df.empty:
                df = df.fillna('')
                tables.append({
                    'source': f'Sheet: {sheet_name}',
                    'data': df.to_dict('records'),
                    'headers': df.columns.tolist(),
                    'merged': None,
                    'method': 'pandas'
                })
    except Exception as e:
        print(f"Error extracting from Excel: {e}")
    if not tables:
        # fallback: convert excel to html and parse
        try:
            html_out = file_path + ".html"
            pd.read_excel(file_path).to_html(html_out)
            html_tables = extract_tables_from_html(html_out)
            for t in html_tables:
                t['source'] = f'Excel->HTML Fallback - {t["source"]}'
                t['method'] = "excel2html"
            tables += html_tables
            os.remove(html_out)
        except Exception as e:
            print(f"Excel fallback error: {e}")
    return tables

def extract_tables_from_csv(file_path):
    """Extract tables from CSV, fallback to text parsing if none found."""
    tables = []
    try:
        df = pd.read_csv(file_path)
        if not df.empty:
            df = df.fillna('')
            tables.append({
                'source': 'CSV File',
                'data': df.to_dict('records'),
                'headers': df.columns.tolist(),
                'merged': None,
                'method': 'pandas'
            })
    except Exception as e:
        print(f"Error extracting from CSV: {e}")
    if not tables:
        # fallback: parse as text, split on common delimiters
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
            if not lines:
                return []
            header = re.split(r'[,\t;|]', lines[0].strip())
            data = [dict(zip(header, re.split(r'[,\t;|]', row.strip()) + [""] * (len(header) - len(re.split(r'[,\t;|]', row.strip()))))) for row in lines[1:]]
            tables.append({
                'source': 'CSV (Text Fallback)',
                'data': data,
                'headers': header,
                'merged': None,
                'method': 'text'
            })
        except Exception as e:
            print(f"CSV fallback error: {e}")
    return tables

def extract_tables_from_html(file_path):
    """Extract tables from HTML using BeautifulSoup or pandas fallback."""
    tables = []
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
            html_tables = soup.find_all('table')
            for table_num, table in enumerate(html_tables, 1):
                rows = table.find_all('tr')
                if rows:
                    headers = [th.get_text(strip=True) for th in rows[0].find_all(['th', 'td'])]
                    data = []
                    for row in rows[1:]:
                        cells = row.find_all(['td', 'th'])
                        row_data = [cell.get_text(strip=True) for cell in cells]
                        if len(row_data) == len(headers):
                            data.append(dict(zip(headers, row_data)))
                    tables.append({
                        'source': f'HTML Table {table_num}',
                        'data': data,
                        'headers': headers,
                        'merged': None,
                        'method': 'bs4'
                    })
    except Exception as e:
        print(f"Error extracting from HTML: {e}")
    if not tables:
        # fallback: try pd.read_html
        try:
            dfs = pd.read_html(file_path)
            for idx, df in enumerate(dfs, 1):
                df = df.fillna('')
                tables.append({
                    'source': f'HTML Table (pandas) {idx}',
                    'data': df.to_dict('records'),
                    'headers': df.columns.tolist(),
                    'merged': None,
                    'method': 'pandas'
                })
        except Exception as e:
            print(f"HTML fallback error: {e}")
    return tables

def extract_tables_from_file(file_path, filename):
    """Universal entrypoint: extract tables from any supported file type, with fallback logic."""
    file_ext = filename.lower().split('.')[-1]
    tables = []

    # Helper: fallback to Azure if nothing found
    def try_azure_fallback(path, label_prefix=""):
        print("🟡 Trying Azure fallback...")
        azure_tables = extract_tables_with_azure_layout(path)
        for t in azure_tables:
            t['source'] = f"{label_prefix}{t['source']}"
            t['method'] = 'Azure'
        return azure_tables

    # ----- PDF -----
    if file_ext == 'pdf':
        print(f"📄 Trying pdfplumber on: {file_path}")
        tables = extract_tables_from_pdf(file_path)
        if not tables:
            tables = try_azure_fallback(file_path, label_prefix="PDF Fallback - ")

    # ----- PPTX / PPT -----
    elif file_ext in ['pptx', 'ppt']:
        print(f"📊 Trying python-pptx on: {file_path}")
        tables = extract_tables_from_pptx(file_path)
        if not tables:
            pdf_out = file_path + ".pdf"
            convert_pptx_to_pdf(file_path, pdf_out)

            from pdf2image import convert_from_path
            poppler_path = r"C:\Users\rohit\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin"
            page_images = convert_from_path(pdf_out, poppler_path=poppler_path)

            for idx, img in enumerate(page_images, 1):
                img_path = f"{file_path}.page{idx}.png"
                img.save(img_path)
                print(f"✅ Saved image: {img_path}")
                tables += try_azure_fallback(img_path, label_prefix=f"PPT Page {idx} - ")
                os.remove(img_path)
            os.remove(pdf_out)

    # ----- Images -----
    elif file_ext in ['jpg', 'jpeg', 'png', 'webp', 'bmp', 'tiff']:
        tables = extract_tables_with_azure_layout(file_path)

    # ----- DOCX -----
    elif file_ext == 'docx':
        tables = extract_tables_from_docx(file_path)

    # ----- Excel -----
    elif file_ext in ['xlsx', 'xls']:
        tables = extract_tables_from_excel(file_path)

    # ----- CSV -----
    elif file_ext == 'csv':
        tables = extract_tables_from_csv(file_path)

    # ----- HTML -----
    elif file_ext in ['html', 'htm']:
        tables = extract_tables_from_html(file_path)

    # ----- Final fallback -----
    if not tables:
        print("⚠️ No tables found from primary methods, trying Azure as last resort...")
        tables = extract_tables_with_azure_layout(file_path)

    print("✅ Final tables to render:", tables)
    return tables

def extract_tables_with_azure_layout(file_path):
    """Extract tables using Azure Form Recognizer Layout model."""
    try:
        with open(file_path, "rb") as f:
            poller = azure_client.begin_analyze_document("prebuilt-layout", document=f)
            result = poller.result()
        
        tables = []
        for t_idx, table in enumerate(result.tables, 1):
            table_data = []
            for row_idx in range(table.row_count):
                row_data = []
                for col_idx in range(table.column_count):
                    # Find cell at this position
                    cell = next(
                        (c for c in table.cells if c.row_index == row_idx and c.column_index == col_idx),
                        None
                    )
                    if cell:
                        row_data.append({
                            "value": cell.content,
                            "colspan": cell.column_span if cell.column_span else 1,
                            "rowspan": cell.row_span if cell.row_span else 1
                        })
                    else:
                        row_data.append({
                            "value": "",
                            "colspan": 1,
                            "rowspan": 1
                        })
                table_data.append(row_data)

            # Optionally extract headers (e.g. first row as header labels)
            headers = [cell["value"] for cell in table_data[0]] if table_data else []

            tables.append({
                'source': f'Azure Layout Table {t_idx}',
                'headers': headers,
                'data': None,
                'merged': table_data,
                'method': 'Azure Layout'
            })

        return tables
    except Exception as e:
        print(f"Azure extraction error: {e}")
        return []

# -------------------- HTML Template --------------------
HTML_TEMPLATE = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Table Extractor</title>
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
            background-color: #f5f5f5;
        }
        .container {
            background-color: white;
            padding: 30px;
            border-radius: 10px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }
        h1 {
            color: #333;
            text-align: center;
            margin-bottom: 30px;
        }
        .upload-form {
            border: 2px dashed #ddd;
            padding: 40px;
            text-align: center;
            margin-bottom: 30px;
            border-radius: 10px;
            background-color: #fafafa;
        }
        input[type="file"] {
            margin: 20px 0;
            padding: 10px;
            font-size: 16px;
        }
        button {
            background-color: #007bff;
            color: white;
            padding: 12px 30px;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            font-size: 16px;
        }
        .results {
            margin-top: 30px;
        }
        table {
            width: 100%;
            border-collapse: collapse;
            background-color: white;
            border: 1.5px solid #222;
        }
        th, td {
            padding: 12px;
            text-align: left;
            border: 1.5px solid #222;
        }
        th {
            background-color: #f8f9fa;
            font-weight: bold;
        }
        tr:hover {
            background-color: #f5f5f5;
        }
        .no-tables {
            text-align: center;
            color: #666;
            font-style: italic;
            padding: 40px;
        }
        .method-label {
            margin-top: 5px;
            color: #666;
            font-size: 0.95em;
        }
    </style>
</head>
<body>
<div class="container">
    <h1>🔍 Universal Table Extractor</h1>
    <form method="POST" enctype="multipart/form-data" class="upload-form">
        <h3>📁 Upload your file to extract tables</h3>
        <p>Support for PDF, Word, PowerPoint, Excel, CSV, HTML, and Images (including WebP)</p>
        <input type="file" name="file" accept=".pdf,.docx,.doc,.pptx,.ppt,.xlsx,.xls,.csv,.txt,.html,.htm,.jpg,.jpeg,.png,.gif,.bmp,.tiff,.webp" required>
        <br>
        <button type="submit">Extract Tables</button>
        <div class="supported-formats">
            <strong>Supported formats:</strong> PDF, DOCX, DOC, PPTX, PPT, XLSX, XLS, CSV, HTML, HTM, JPG, JPEG, PNG, GIF, BMP, TIFF, WEBP
        </div>
    </form>

    {% if tables %}
    <div class="results">
        <h2>📊 Extracted Tables ({{ tables|length }} found)</h2>
        {% for table in tables %}
        <div class="table-container">
            <h3>{{ table.source }}</h3>
            <div class="method-label">Extracted by: <b>{{ table.method }}</b></div>
            {% if table.data or table.merged %}
            <table id="table-{{ loop.index }}">
                <thead>
    {% if table.headers and table.headers[0] is mapping %}
    <!-- Handle merged header cells from Gemini -->
    {% for row in table.headers %}
    <tr>
        {% for cell in row %}
        <th colspan="{{ cell.colspan|default(1) }}" rowspan="{{ cell.rowspan|default(1) }}">
            {{ cell.value }}
        </th>
        {% endfor %}
    </tr>
    {% endfor %}
    {% else %}
    <tr>
        {% for header in table.headers %}
        <th>{{ header }}</th>
        {% endfor %}
    </tr>
    {% endif %}
</thead>

                <tbody>
                    {% if table.data %}
                        {% for row in table.data %}
                        <tr>
                            {% for header in table.headers %}
                            <td>{{ row.get(header, '') or '' }}</td>
                            {% endfor %}
                        </tr>
                        {% endfor %}
                    {% elif table.merged %}
                        {% for row in table.merged %}
                        <tr>
                            {% for cell in row %}
                            <td colspan="{{ cell.colspan|default(1) }}" rowspan="{{ cell.rowspan|default(1) }}">{{ cell.value }}</td>
                            {% endfor %}
                        </tr>
                        {% endfor %}
                    {% endif %}
                </tbody>
            </table>
            {% else %}
            <div class="no-tables">No data found in this table</div>
            {% endif %}
        </div>
        {% endfor %}
    </div>
    {% endif %}

    {% if tables is defined and tables|length == 0 %}
    <div class="no-tables">
        <h3>🤷‍♂️ No tables found in the uploaded file</h3>
        <p>The file was processed successfully, but no tables were detected.</p>
    </div>
    {% endif %}
</div>
</body>
</html>
'''

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    """Main route for uploading and extracting tables from files."""
    if request.method == 'POST':
        if 'file' not in request.files:
            flash('No file selected', 'error')
            return redirect(request.url)
        file = request.files['file']
        if file.filename == '':
            flash('No file selected', 'error')
            return redirect(request.url)
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            try:
                tables = extract_tables_from_file(file_path, filename)
                os.remove(file_path)
                if tables:
                    flash(f'Successfully extracted {len(tables)} table(s) from {filename}', 'success')
                else:
                    flash(f'No tables found in {filename}', 'error')
                return render_template_string(HTML_TEMPLATE, tables=tables)
            except Exception as e:
                if os.path.exists(file_path):
                    os.remove(file_path)
                flash(f'Error processing file: {str(e)}', 'error')
                return render_template_string(HTML_TEMPLATE)
        else:
            flash('File type not supported', 'error')
            return redirect(request.url)
    return render_template_string(HTML_TEMPLATE)

@app.route('/api/extract', methods=['POST'])
def api_extract():
    """API endpoint for extracting tables (JSON response)."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    if not allowed_file(file.filename):
        return jsonify({'error': 'File type not supported'}), 400
    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)
    try:
        tables = extract_tables_from_file(file_path, filename)
        os.remove(file_path)
        return jsonify({
            'success': True,
            'filename': filename,
            'tables_found': len(tables),
            'tables': tables
        })
    except Exception as e:
        if os.path.exists(file_path):
            os.remove(file_path)
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("🚀 Starting Table Extractor Application with Donut + Gemini + fallback for all file types...")
    app.run(debug=True, host='0.0.0.0', port=5000)